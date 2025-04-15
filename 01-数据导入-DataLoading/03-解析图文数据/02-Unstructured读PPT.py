# 需要安装LibreOffice

# import os
# os.environ["UNSTRUCTURED_LIBREOFFICE_PATH"] = r"C:\Program Files\LibreOffice\program\soffice.exe"
# os.environ["UNSTRUCTURED_LIBREOFFICE_ENCODING"] = "gbk"  # 添加这行解决编码问题

# from unstructured.partition.ppt import partition_ppt

# # ================================================使用原生unstructured解析pptx文件=====================================
# # 解析 PPT 文件
# ppt_elements = partition_ppt(filename="90-文档-Data/黑悟空/探索马尔代夫的魅力.pptx")
# print("PPT 内容：")
# for element in ppt_elements:
#     print(element.text)
    
# ================================================使用langchain_core将unstructured解析的pptx文件转换为Documents格式=====================================
from langchain_core.documents import Document
from pptx import Presentation
import os
ppt_path = os.path.abspath("90-文档-Data/黑悟空/探索马尔代夫的魅力.pptx")

# 使用python-pptx直接解析PPT
prs = Presentation(ppt_path)
text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_runs.append(shape.text)

# 转换为 Documents 数据结构
# 转换为Documents格式
documents = [
    Document(page_content=text, 
            metadata={"source": ppt_path})
    for text in text_runs
]

# 输出转换后的 Documents
print(documents[0:])


